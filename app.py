import streamlit as st
import pandas as pd
from pandasai import SmartDataframe
from pandasai import SmartDatalake
from pandasai.llm import OpenAI
import os
from dotenv import load_dotenv
import openpyxl
from openpyxl import Workbook
from pandasai.helpers.openai_info import get_openai_callback
import time
import os
from openai import OpenAI as OpenAIClient
from pandasai.connectors import PandasConnector
import ast
from pptx import Presentation
import json
from datetime import datetime
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pdf2image import convert_from_path
import fitz  # PyMuPDF
import io
from PIL import Image
import tempfile
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64
from io import BytesIO
import numpy as np
import cv2
import shutil  # For directory cleanup
from utils import data_dictt, FAQs, mark_D, dril_ppt, drill_p_m, consolidation_prompt, consolidatin_prompt_main
from database import DocumentDatabase
import hashlib

# Imports from filtered.py
from temp_filtered import * 
from infographic_functions import *
# -----------------------
# At top of app.py (right after imports)
for key, val in {
    "chat_history": [],
    "should_run_filtered": False,
    "should_show_selection": False,
    "last_query": None,
    "filtered_results": None,   # whatever filtered.py returns
    "selected_prompts": None,   # user's filtered selection
    "filtered_answers": None,   # answered filtered queries
    "show_unfiltered": False,
    "final_answer": None,
}.items():
    if key not in st.session_state:
        st.session_state[key] = val


UPLOAD_DIR = "uploaded_files"
os.makedirs(UPLOAD_DIR, exist_ok=True)
# -----------------------

# Keep only these lines after import
st.cache_data.clear()
st.cache_resource.clear()
st.session_state.should_query = False
load_dotenv()

# Initialize database
db = DocumentDatabase()

# Initialize session state for storing prompts and feedback
if "previous_prompts" not in st.session_state:
    st.session_state.previous_prompts = []
if 'chat_history' not in st.session_state:
    st.session_state['chat_history'] = []
if "should_query" not in st.session_state:
    st.session_state.should_query = False
if "use_tesseract" not in st.session_state:
    st.session_state.use_tesseract = False
if "debug_mode" not in st.session_state:
    st.session_state.debug_mode = False
if "current_query" not in st.session_state:
    st.session_state.current_query = None
# Recently Updated.
if 'triangulation' not in st.session_state:
    st.session_state['triangulation'] = None
if 'main_answer_df' not in st.session_state:
    st.session_state['main_answer_df'] = None
if "database_previous_prompts" not in st.session_state:
    st.session_state.database_previous_prompts = []
if "last_generated_sql" not in st.session_state:
    st.session_state['last_generated_sql'] = ''
if "query_sql_map" not in st.session_state:
    st.session_state['query_sql_map'] = {}
if "selected_sql_for_edit" not in st.session_state:
    st.session_state['selected_sql_for_edit'] = ''
# if "current_query_token_usage" not in st.session_state:
#     st.session_state.current_query_token_usage = 0
# token_usage_placeholder = st.empty()

def add_token_usage(token_count: int, query_text: str, usage_type: str = "LLM Query"):
    if "current_query_token_usage" not in st.session_state:
        st.session_state.current_query_token_usage = 0
    if "token_usage_history" not in st.session_state:
        st.session_state.token_usage_history = []

    st.session_state.current_query_token_usage += token_count
    st.session_state.token_usage_history.append({
        "query": query_text,
        "tokens": token_count,
        "type": usage_type
    })

token_placeholder = st.empty()

def display_token_usage():
    token_count = st.session_state.get("current_query_token_usage", 0)
    token_placeholder.markdown(
        f"""
        <div style="
            position: fixed;
            bottom: 10px;
            right: 10px;
            background-color: #f1f1f1;
            padding: 8px 12px;
            border-radius: 8px;
            box-shadow: 0 0 8px rgba(0,0,0,0.1);
            font-size: 14px;
            color: #333;
            z-index: 1000;
        ">
            Total tokens used: <b>{token_count}</b>
        </div>
        """, unsafe_allow_html=True
    )

# Call this once per run, after token updates (or at end of main script)
# display_token_usage()


# Initialize lists for storing processed documents
pdf_texts = []
ppt_texts = []
dataframes = []
extracted_images = []

unique_values_dict = {}
data_dict = data_dictt

# Define FAQs categorized by function
frequently_asked_prompts_by_function = FAQs

# Page title
st.title("Priorise Insights Engine")

def calculate_file_hash(file_content):
    """Calculate SHA-256 hash of file content"""
    return hashlib.sha256(file_content).hexdigest()

# Add custom CSS to hide file uploader size limit text
st.markdown(mark_D, unsafe_allow_html=True)
# Near the top of your file, after the title
tab1, tab2, tab3 = st.tabs(["Main", "SQL", "Chat History"])

# Initialize OpenAI client
llm = OpenAI(api_token=os.getenv("OPEN_API_KEY"), temperature=0, seed=42, model="gpt-4o")
client = OpenAIClient(api_key=os.getenv("OPEN_API_KEY"))

# ----------------------------------------------------------------------------------------------
# Here Writing Engine
from temp_filtered import QueryProcessor, KPISelector, SmartQueryEngine, schema_json, kpi_dict
client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
# y_processor = QueryProcessor(openai_api_key=os.getenv("OPENAI_API_KEY"))
kpi_selector = KPISelector()
DB_CONFIG = {
    "host": os.getenv("DB_HOST"),
    "database": os.getenv("DB_NAME"),
    "user": os.getenv("DB_USER"),
    "password": os.getenv("DB_PASSWORD"),
    "port": os.getenv("DB_PORT")
    }

if 'engine' not in st.session_state:
    st.session_state.engine = SmartQueryEngine(
        client=client,
        schema_json=schema_json,
        db_config=DB_CONFIG,
        kpi_dict=kpi_dict
    )

if 'query_processor' not in st.session_state:
    st.session_state.query_processor = QueryProcessor(openai_api_key=os.getenv("OPENAI_API_KEY"))
# -----------------------------------------------------------------------------------------------

# Add after imports
if st.session_state.debug_mode:
    def check_poppler_installation():
        try:
            # Check if poppler is in PATH
            paths = os.environ["PATH"].split(os.pathsep)
            st.write("Checking PATH for Poppler:")
            poppler_found = False
            for path in paths:
                if "poppler" in path.lower():
                    st.write(f"Found Poppler in PATH: {path}")
                    poppler_found = True
            
            if not poppler_found:
                st.warning("Poppler not found in PATH")
            
            # Try to execute pdftoppm (a Poppler utility)
            import subprocess
            try:
                subprocess.run(['pdftoppm', '-v'], capture_output=True)
                st.success("Poppler utilities are accessible!")
            except FileNotFoundError:
                st.error("Poppler utilities not found in system PATH")
        except Exception as e:
            st.error(f"Error checking Poppler installation: {str(e)}")
    
    # Add a debug section in the sidebar
    with st.sidebar:
        with st.expander("Debug Poppler Installation"):
            check_poppler_installation()

# In tab1 (Main), put all your existing main content
with tab1:
    # Sidebar content
    with st.sidebar:
        st.image("Logo.jpg", width=280)
        selected_function = st.selectbox("Choose a function", list(frequently_asked_prompts_by_function.keys()))
        with st.expander("Frequently Asked Prompts"):
            st.write(f"Here are some frequently asked questions for the {selected_function} function:")
            for i, prompt in enumerate(frequently_asked_prompts_by_function[selected_function], 1):
                st.write(f"{i}. {prompt}")
        st.subheader("Upload Your Datasets")

        # Updated
        # Initialize storage if not present
        uploaded_files = st.file_uploader('Choose CSV/Excel/PPT/PDF files', type=["xlsx", "csv", "pptx", "pdf"], accept_multiple_files=True)
        if uploaded_files:
            st.write(f"{len(uploaded_files)} file(s) uploaded.")

        st.write("---")
        
        # Dynamic Column Names Display
        if uploaded_files:
            for file in uploaded_files:
                if file.name.endswith('.csv'):
                    df = pd.read_csv(file)
                    
                    with st.expander(f"Data Dictionary"):
                        for col in df.columns:
                            st.write(f'- {col} - {data_dict[col]}', unsafe_allow_html=True)
                elif file.name.endswith('.xlsx'):
                    df = pd.read_excel(file)
                    with st.expander(f"Data Dictionary"):
                        for col in df.columns:
                            st.write(f'- {col} - {data_dict[col]}', unsafe_allow_html=True)
        # st.session_state.debug_mode = st.checkbox("Enable Debug Mode", value=st.session_state.debug_mode)
        # st.session_state.use_tesseract = st.checkbox("Use Tesseract OCR", value=st.session_state.use_tesseract)
        # Assume this dictionary saved in session state:
    # Main content
    # st.write("---")
    st.subheader("How can I help you?")
    temp_query = st.text_input("Enter your question:")

    # Ask button and save query to chat history
    col1, col2, col3, empty_col = st.columns([0.9, 1.6, 1, 1.8])
    with col1:
        if st.button("Sql Query", key="ask_btn"):
            if temp_query:
                # Reset token usage for the new query
                st.session_state.current_query_token_usage = 0

                st.session_state.should_query = True
                st.session_state['chat_history'].append({"query": temp_query, "response": None})
                st.session_state.current_query = temp_query  # Store the current query
                # Updated
                st.session_state.should_run_filtered = True
                st.session_state.last_query = temp_query
                st.session_state['main_answer_df'] = None
                st.session_state.database_previous_prompts.append(temp_query)
    
    with col2:
        if st.button("Query Knowledge Base", key="kb_btn"):
            if temp_query:
                st.session_state.current_query_token_usage = 0
                st.session_state.should_query_kb = True
                st.session_state.current_query = temp_query
    
    with col3:
        if st.button("Triangulation", key="Tri_btn"):
            if temp_query:
                st.session_state.current_query_token_usage = 0
                st.session_state.should_query = True
                st.session_state.current_query = temp_query
                st.session_state.last_query = temp_query
                st.session_state.triangulation = True


    if temp_query:
        st.markdown(f"<p style='font-size:20px; font-weight:bold;'>{temp_query}</p>", unsafe_allow_html=True)
    else:
        st.write("No query entered yet.")

    if st.session_state.get('should_query_kb', False):
        st.session_state['main_answer_df'] = None #Setting it to none because is should not display the main answer Dataframe.
        # ---- Start your inline logic here ----
        kb_query = st.session_state.current_query
        # st.write(f"Querying Knowledge Base for: {kb_query}")
        st.markdown("**Query Knowledge Base:**")
        if not uploaded_files:
                st.warning("Please Upload files for Query Knowledge Base analysis!")
        # Add your normal app.py fallback answer logic here
        EXCEL_FILE_PATH = "callback_data.xlsx"
        user_defined_path = os.getcwd()
        udp = os.path.join(user_defined_path, "exports", "charts")

        # Function to parse callback output
        def parse_callback_output(cb) -> dict:
            output_lines = str(cb).split("\n")
            parsed_output = {
                "total_tokens": int(output_lines[0].split(":")[1].strip()),
                "prompt_tokens": int(output_lines[1].split(":")[1].strip()),
                "completion_tokens": int(output_lines[2].split(":")[1].strip()),
                "total_cost": float(output_lines[3].split(":")[1].strip().replace("$", "").strip()),
            }
            return parsed_output

        # Function to append data to Excel
        def append_to_excel(file_path: str, data: dict):
            if not os.path.exists(file_path):
                wb = Workbook()
                ws = wb.active
                ws.append(["Total Tokens", "Prompt Tokens", "Completion Tokens", "Total Cost (USD)"])
                wb.save(file_path)
            wb = openpyxl.load_workbook(file_path)
            ws = wb.active
            ws.append([data["total_tokens"], data["prompt_tokens"], data["completion_tokens"], data["total_cost"]])
            wb.save(file_path)


        # Add these new functions after the imports
        def extract_images_from_pdf(pdf_file):
            """Extract images and charts from PDF using PyMuPDF with focus on chart detection"""
            images = []
            temp_file = None
            doc = None
            
            try:
                # Create a temporary file to save the uploaded PDF
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                temp_file_path = temp_file.name
                temp_file.write(pdf_file.read())
                temp_file.close()
                
                # Open the PDF file
                doc = fitz.open(temp_file_path)
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    
                    # Method 1: Standard image extraction
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # Convert to PIL Image
                        image = Image.open(io.BytesIO(image_bytes))
                        if is_valid_image_size(image):
                            images.append({
                                'page': page_num + 1,
                                'index': img_index,
                                'image': image,
                                'format': base_image["ext"],
                                'method': 'direct_extraction',
                                'width': image.width,
                                'height': image.height
                            })
                    
                    # Method 2: Extract potential chart regions by analyzing page content
                    # Get all drawings on the page
                    paths = page.get_drawings()
                    if paths:
                        # Extract regions with significant drawing content
                        regions = analyze_drawing_regions(paths)
                        for region_index, region in enumerate(regions):
                            try:
                                # Render the region with higher resolution
                                zoom = 2  # Increase if needed for better quality
                                mat = fitz.Matrix(zoom, zoom)
                                pix = page.get_pixmap(matrix=mat, clip=region)
                                
                                # Convert to PIL Image
                                img_data = Image.frombytes("RGB", [pix.width, pix.height], pix.tobytes())
                                
                                if is_valid_image_size(img_data):
                                    images.append({
                                        'page': page_num + 1,
                                        'index': f'chart_{region_index}',
                                        'image': img_data,
                                        'format': 'png',
                                        'method': 'region_extraction',
                                        'width': img_data.width,
                                        'height': img_data.height,
                                        'bbox': [region.x0, region.y0, region.x1, region.y1]
                                    })
                            except Exception as e:
                                print(f"Error extracting region: {e}")
                
                return images
                
            except Exception as e:
                st.error(f"Error extracting images from PDF: {e}")
                return []
            
            finally:
                if doc:
                    doc.close()
                if temp_file and os.path.exists(temp_file.name):
                    try:
                        os.unlink(temp_file.name)
                    except Exception as e:
                        print(f"Error removing temporary file: {e}")

        def is_valid_image_size(image):
            """Check if image meets minimum size requirements"""
            MIN_WIDTH = 100
            MIN_HEIGHT = 100
            MAX_DIMENSION = 4000  # Prevent extremely large images
            
            return (MIN_WIDTH <= image.width <= MAX_DIMENSION and 
                    MIN_HEIGHT <= image.height <= MAX_DIMENSION)

        def analyze_drawing_regions(paths):
            """Analyze drawing paths to identify potential chart regions"""
            regions = []
            
            if not paths:
                return regions
            
            # Group nearby paths that might form a chart
            current_group = []
            
            for path in paths:
                rect = path['rect']  # Get the bounding rectangle
                
                # Skip tiny drawings
                if rect.width < 20 or rect.height < 20:
                    continue
                
                if current_group:
                    # Check if this path is close to the current group
                    last_rect = current_group[-1]['rect']
                    if (abs(rect.x0 - last_rect.x1) < 50 and 
                        abs(rect.y0 - last_rect.y1) < 50):
                        current_group.append(path)
            else:
                    current_group.append(path)
            
            # Don't forget the last group
            if len(current_group) >= 3:
                merged_rect = merge_rects([p['rect'] for p in current_group])
                if merged_rect.width >= 100 and merged_rect.height >= 100:
                    regions.append(merged_rect)
            
            return regions

        def merge_rects(rects):
            """Merge multiple rectangles into one bounding rectangle"""
            if not rects:
                return None
            
            x0 = min(rect.x0 for rect in rects)
            y0 = min(rect.y0 for rect in rects)
            x1 = max(rect.x1 for rect in rects)
            y1 = max(rect.y1 for rect in rects)
            
            # Add some padding
            padding = 10
            return fitz.Rect(x0 - padding, y0 - padding, x1 + padding, y1 + padding)

        def extract_images_from_pptx(pptx_file):
            """Extract images from PowerPoint file"""
            images = []
            try:
                prs = Presentation(pptx_file)
                
                for slide_num, slide in enumerate(prs.slides):
                    for shape_num, shape in enumerate(slide.shapes):
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_stream = io.BytesIO(shape.image.blob)
                            image = Image.open(image_stream)
                            images.append({
                                'slide': slide_num + 1,
                                'index': shape_num,
                                'image': image,
                                'format': image.format.lower()
                            })
                return images
            except Exception as e:
                st.error(f"Error extracting images from PPTX: {e}")
                return []

        def convert_pptx_to_pdf(input_dir, output_dir):
            """Convert PPTX to PDF using platform-specific methods"""
            try:
                if os.name == 'nt':  # Windows
                    # Import pythoncom only for Windows
                    import pythoncom
                    # Initialize COM for this thread
                    pythoncom.CoInitialize()
                    try:
                        from pptxtopdf import convert
                        convert(input_dir, output_dir)
                    finally:
                        # Uninitialize COM
                        pythoncom.CoUninitialize()
                else:  # Linux/Unix
                    # Check if libreoffice is installed
                    try:
                        import subprocess
                        subprocess.run(['libreoffice', '--version'], capture_output=True, check=True)
                    except (subprocess.CalledProcessError, FileNotFoundError):
                        raise RuntimeError("LibreOffice is not installed. Please install it using: sudo apt-get install libreoffice")
                    
                    # Get the input file path
                    input_file = os.path.join(input_dir, os.listdir(input_dir)[0])  # Get the first file in the directory
                    
                    # Convert using libreoffice with explicit file path
                    subprocess.run([
                        'libreoffice',
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', output_dir,
                        input_file
                    ], check=True)
                    
                    # Add debug logging
                    if st.session_state.debug_mode:
                        st.write(f"Input file: {input_file}")
                        st.write(f"Output directory: {output_dir}")
                        st.write(f"Files in output directory: {os.listdir(output_dir)}")
                    
            except Exception as e:
                raise RuntimeError(f"Error converting PPTX to PDF: {str(e)}")

        def process_pptx_with_gpt4v(pptx_file):
            """Process PPTX using GPT-4 Vision by first converting to PDF then processing like a PDF"""
            # Calculate file hash
            file_content = pptx_file.read()
            file_hash = calculate_file_hash(file_content)
            
            # Check if file exists in database with same hash
            doc_data = db.get_document(pptx_file.name)
            if doc_data and doc_data["file_hash"] == file_hash:
                if st.session_state.debug_mode:
                    st.write(f"Using cached data from database for PPT: {pptx_file.name}")
                return doc_data["content"]
            
            # Reset file pointer for processing
            pptx_file.seek(0)
            
            # Create temporary directories
            temp_input_dir = tempfile.mkdtemp()
            temp_output_dir = tempfile.mkdtemp()
            
            try:
                # Save uploaded file to temp directory
                temp_input_path = os.path.join(temp_input_dir, pptx_file.name)
                with open(temp_input_path, 'wb') as f:
                    f.write(file_content)
                
                # Convert PPTX to PDF using platform-specific method
                convert_pptx_to_pdf(temp_input_dir, temp_output_dir)
                
                # Find the converted PDF file
                pdf_filename = os.path.splitext(pptx_file.name)[0] + '.pdf'
                pdf_path = os.path.join(temp_output_dir, pdf_filename)
                
                if not os.path.exists(pdf_path):
                    raise FileNotFoundError(f"PDF file not found at expected location: {pdf_path}")
                
                if st.session_state.debug_mode:
                    st.write(f"Successfully converted PPTX to PDF: {pdf_path}")
                
                # Process the PDF using our existing PDF processing logic
                with open(pdf_path, 'rb') as pdf_file:
                    # Initialize OpenAI client
                    # client = OpenAIClient()
                    
                    # Convert PDF pages to images
                    images = convert_from_path(pdf_path)
                    page_texts = []
                    
                    for i, image in enumerate(images):
                        # Convert image to base64
                        buffered = BytesIO()
                        image.save(buffered, format="PNG")
                        base64_image = base64.b64encode(buffered.getvalue()).decode()
                        
                        # Extract text using GPT-4V
                        response = client.chat.completions.create(
                            model="gpt-4o",
                            messages=[
                                {
                                    "role": "user",
                                    "content": [
                                        {
                                            "type": "text",
                                            "text": "Extract all content from this slide. Include text, describe any images, charts, or diagrams, and maintain the formatting and structure."
                                        },
                                        {
                                            "type": "image_url",
                                            "image_url": {
                                                "url": f"data:image/png;base64,{base64_image}",
                                                "detail": "high"
                                            }
                                        }
                                    ]
                                }
                            ],
                            max_tokens=1000
                        )
                        
                        extracted_text = response.choices[0].message.content
                        tokens_used = response.usage.total_tokens
                        add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                        page_texts.append({
                            "slide_number": i + 1,
                            "text": extracted_text,
                            "char_count": len(extracted_text),
                            "word_count": len(extracted_text.split())
                        })
                        
                        if st.session_state.debug_mode:
                            with st.expander(f"Debug: Extracted Content Slide {i+1}", expanded=False):
                                st.write("#### Extracted Content")
                                st.text_area("Text Content", extracted_text, height=300)
                                st.write(f"Characters extracted: {len(extracted_text)}")
                                st.write(f"Words extracted: {len(extracted_text.split())}")
                    
                    # Combine all slide texts
                    combined_text = "\n\n".join([page["text"] for page in page_texts])
                    result = {
                        "text": combined_text,
                        "slides/pages": page_texts,
                        "processed_date": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    }
                    
                    # Save to database
                    db.save_document(
                        filename=pptx_file.name,
                        file_type="pptx",
                        content=result,
                        file_hash=file_hash
                    )
                    
                    return result
                    
            except Exception as e:
                if st.session_state.debug_mode:
                    st.error(f"Error during PDF conversion/processing: {str(e)}")
                raise e
                
            finally:
                # Clean up temporary directories
                try:
                    if os.path.exists(temp_input_dir):
                        shutil.rmtree(temp_input_dir)
                    if os.path.exists(temp_output_dir):
                        shutil.rmtree(temp_output_dir)
                except Exception as e:
                    if st.session_state.debug_mode:
                        st.error(f"Error cleaning up temporary files: {str(e)}")

        # Add new function for GPT-4 Vision extraction
        def extract_text_with_gpt4v(pdf_file):
            """Extract text from PDF using GPT-4 Vision"""
            try:
                # Create a temporary file to save the uploaded PDF
                temp_file = None
                text_results = []
                
                try:
                    # Save uploaded file to temporary file
                    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                    temp_file.write(pdf_file.read())
                    temp_file.close()
                    
                    try:
                        # Convert PDF pages to images using pdf2image
                        images = convert_from_path(temp_file.name)
                        
                        # Process each page
                        for i, image in enumerate(images):
                            if st.session_state.debug_mode:
                                with st.expander(f"Debug: GPT-4V Processing Page {i+1}", expanded=False):
                                    st.write("---")
                                    st.write(f"### Processing Page {i+1}")
                                    st.image(image, caption=f"Page {i+1}", use_column_width=True)
                            
                            # Convert PIL Image to bytes
                            img_byte_arr = io.BytesIO()
                            image.save(img_byte_arr, format='PNG')
                            img_byte_arr = img_byte_arr.getvalue()
                            
                            # Encode image to base64
                            base64_image = base64.b64encode(img_byte_arr).decode('utf-8')
                            
                            # Create GPT-4V prompt with correct image URL format
                            response = client.chat.completions.create(
                                model="gpt-4o",
                                messages=[
                                    {
                                        "role": "user",
                                        "content": [
                                            {
                                                "type": "text",
                                                "text": "Extract all text content from this image. Preserve formatting and structure. Include any relevant tables, charts, or diagrams descriptions."
                                            },
                                            {
                                                "type": "image_url",
                                                "image_url": {
                                                    "url": f"data:image/png;base64,{base64_image}",
                                                    "detail": "high"
                                                }
                                            }
                                        ]
                                    }
                                ],
                                max_tokens=1000
                            )
                            
                            extracted_text = response.choices[0].message.content
                            tokens_used = response.usage.total_tokens
                            add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")

                            if st.session_state.debug_mode:
                                with st.expander(f"Debug: Extracted Content Page {i+1}", expanded=False):
                                    st.write("#### Extracted Content")
                                    st.text_area("Text Content", extracted_text, height=300)
                                    st.write(f"Characters extracted: {len(extracted_text)}")
                                    st.write(f"Words extracted: {len(extracted_text.split())}")
                            
                            text_results.append(extracted_text)
                        
                        return "\n\n".join(text_results)
                        
                    except Exception as e:
                        if st.session_state.debug_mode:
                            st.error(f"Error processing PDF: {str(e)}")
                        raise e
                        
                finally:
                    # Clean up temporary file
                    if temp_file and os.path.exists(temp_file.name):
                        try:
                            os.unlink(temp_file.name)
                        except Exception as e:
                            if st.session_state.debug_mode:
                                st.error(f"Error removing temporary file: {e}")
                        
            except Exception as e:
                st.error(f"Error in GPT-4V processing: {e}")
                if st.session_state.debug_mode:
                    st.error(f"Detailed error information: {str(e)}")
                return ""

        # Prepare datasets for SmartDatalake and process PDFs immediately
        if uploaded_files:
            dataframes = []
            pdf_texts = []
            ppt_texts = []  # New list for PPT data
            extracted_images = []
            
            for uploaded_file in uploaded_files:
                try:
                    if uploaded_file.name.endswith('.pptx'):
                        # Process PPTX files
                        uploaded_file.seek(0)
                        pptx_result = process_pptx_with_gpt4v(uploaded_file)
                        if pptx_result:
                            ppt_texts.append({
                                "name": uploaded_file.name,
                                "text": pptx_result["text"],
                                "slides": pptx_result["slides/pages"]
                            })
                        
                        # Extract images from PPTX
                        uploaded_file.seek(0)
                        pptx_images = extract_images_from_pptx(uploaded_file)
                        extracted_images.extend([{
                            'source': uploaded_file.name,
                            'type': 'pptx',
                            **img
                        } for img in pptx_images])
                    
                    elif uploaded_file.name.endswith('.pdf'):
                        # Calculate file hash
                        file_content = uploaded_file.read()
                        file_hash = calculate_file_hash(file_content)
                        
                        # Check if file exists in database with same hash
                        doc_data = db.get_document(uploaded_file.name)
                        if doc_data and doc_data["file_hash"] == file_hash:
                            if st.session_state.debug_mode:
                                st.write(f"Using cached data from database for PDF: {uploaded_file.name}")
                            pdf_texts.append({
                                "name": uploaded_file.name,
                                "text": doc_data["content"]["text"],
                                "pages": doc_data["content"]["pages"]
                            })
                        else:
                            # Process new PDF file
                            uploaded_file.seek(0)
                            pdf_reader = PdfReader(uploaded_file)
                            text = ""
                            pages = []
                            for i, page in enumerate(pdf_reader.pages):
                                page_text = page.extract_text() or ""
                                text += page_text
                                pages.append({
                                    "page_number": i + 1,
                                    "text": page_text,
                                    "char_count": len(page_text),
                                    "word_count": len(page_text.split())
                                })
                            
                            result = {
                                "text": text,
                                "pages": pages,
                                "processed_date": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                            }
                            
                            # Save to database
                            db.save_document(
                                filename=uploaded_file.name,
                                file_type="pdf",
                                content=result,
                                file_hash=file_hash
                            )
                            
                            pdf_texts.append({
                                "name": uploaded_file.name,
                                "text": text,
                                "pages": pages
                            })
                    
                    elif uploaded_file.name.endswith('.csv'):
                        df = pd.read_csv(uploaded_file)
                        dataframes.append(df)
                    elif uploaded_file.name.endswith('.xlsx'):
                        df = pd.read_excel(uploaded_file)
                        for col in df.select_dtypes(include=['object']).columns:
                            unique_values_dict[col] = df[col].dropna().unique().tolist()
                        dataframes.append(df)
                    else:
                        st.error(f"Unsupported file format: {uploaded_file.name}")
                        continue
                    
                except Exception as e:
                    st.error(f"Error loading file '{uploaded_file.name}': {e}")
                    continue

            # Store extracted images in session state
            st.session_state['extracted_images'] = extracted_images

            # Process PDFs into chunks
            if pdf_texts:
                combined_text = "\n".join([pdf["text"] for pdf in pdf_texts])
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200,
                    length_function=len
                )
                chunks = text_splitter.split_text(combined_text)
                st.session_state['pdf_chunks'] = chunks

            if dataframes:
                # st.write('dataframes')
                sales_data = SmartDatalake(dataframes, config={"llm": llm, "use_error_correction_framework": True, "data_dict": data_dict})
            else:
                sales_data = None
        else:
            sales_data = None


        # Function to generate up to 3 drill-down questions using the LLM's chat method
        def generate_drill_down_questions(query, dataframes):
            # Check if we have brand metrics data by looking at the dataframe structure
            def is_brand_metrics_df(df):
                required_columns = ['Brand', 'Date', 'TOM', 'Spontaneous_Awareness', 'MOUB', 'Purchase_6M']
                return all(col in df.columns for col in required_columns)
            
            # Determine the type of analysis based on the data structure
            has_brand_metrics = any(is_brand_metrics_df(df) for df in dataframes)
            
            if has_brand_metrics:
                drill_prompt_ppt = dril_ppt.format(query=query)
                response_PPT = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": "You are an AI assistant."},
                            {"role": "user", "content": drill_prompt_ppt}
                        ],
                        temperature=0,
                        max_tokens=512
                    )
                tokens_used = response_PPT.usage.total_tokens
                add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                drill_response_PPT = response_PPT.choices[0].message.content
                clean_response_PPT = drill_response_PPT.replace("```python", "").replace("```", "").strip()
                parsed_questions_PPT= ast.literal_eval(clean_response_PPT)
                return {"type":'PPT',"response":parsed_questions_PPT}
            else:
                
                drill_prompt_main = drill_p_m.format(query=query,
                                                        data_dictt=json.dumps(data_dict, indent=2),
                                                        unique_values_dict=json.dumps({col: unique_values_dict[col] for col in unique_values_dict if col in data_dict}, indent=2))
                response_main = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": "You are an AI assistant."},
                            {"role": "user", "content": drill_prompt_main}
                        ],
                        temperature=0,
                        max_tokens=512
                    )
                tokens_used = response_main.usage.total_tokens
                add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                drill_response_main = response_main.choices[0].message.content
                clean_response_main = drill_response_main.replace("```python", "").replace("```", "").strip()
                parsed_questions_main = ast.literal_eval(clean_response_main)
                return {"type":'main',"response":parsed_questions_main}

        # Function to auto-save chat history
        def auto_save_chat_history():
            """Automatically save the chat history after each interaction"""
            if st.session_state['chat_history']:
                # Create a 'chat_history' directory if it doesn't exist
                os.makedirs('chat_history', exist_ok=True)
                
                # Use a fixed filename with current date
                current_date = datetime.now().strftime("%Y%m%d")
                filename = f"chat_history/chat_history_{current_date}.json"
                
                # Add timestamp to each chat entry that doesn't have one
                chat_history_with_timestamp = []
                for chat in st.session_state['chat_history']:
                    chat_entry = chat.copy()
                    if 'timestamp' not in chat_entry:
                        chat_entry['timestamp'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    chat_history_with_timestamp.append(chat_entry)
                    
                try:
                    with open(filename, 'w', encoding='utf-8') as f:
                        json.dump(chat_history_with_timestamp, f, indent=2, ensure_ascii=False)
                except Exception as e:
                    st.error(f"Error auto-saving chat history: {e}")

        # Add this function to filter data before querying
        def preprocess_data_for_state_query(df):
            if 'Region' in df.columns:
                # Keep all rows except those where MARKET starts with 'ALL INDIA'
                return df[~df['Region'].str.startswith('ALL INDIA', na=False)]
            return df

        # Function to retrieve relevant chunks based on the query with error handling
        def retrieve_relevant_chunks(query, chunks, top_k=3):
            if not chunks:
                st.error("No chunks available to process the query.")
                return []
            prompt = f"""
            Given the query: "{query}", and the following document chunks:
            {json.dumps(chunks[:5], indent=2)} (showing first 5 chunks for brevity),
            identify the top {top_k} most relevant chunks to answer the query.
            Return the indices of these chunks as a list (e.g., [0, 2, 4]).
            """
            try:
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": "You are an AI assistant that identifies relevant text chunks."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0,
                    max_tokens=100
                )
                response_content = response.choices[0].message.content.strip()
                tokens_used = response.usage.total_tokens
                add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                indices = ast.literal_eval(response_content)
                if not isinstance(indices, list):
                    raise ValueError("LLM did not return a valid list of indices")
                return [chunks[i] for i in indices if i < len(chunks)]
            except Exception as e:
                print(f"Error retrieving relevant chunks: {e}")
                return chunks[:top_k]  # Fallback to first top_k chunks

        def process_pdf_query_page_wise(query, pdf_data):
            """Process query by analyzing each page's content individually"""
            page_responses = []
            
            # First, analyze each page individually
            for page in pdf_data['pages']:
                # Handle both page_number and slide_number fields
                page_num = page.get('page_number', page.get('slide_number', 1))
                page_prompt = f"""
                Analyzing page {page_num} of the document.
                Content: {page['text']}
                
                Question: {query}
                
                If this page contains relevant information to answer the question, provide the relevant details.
                If this page does not contain relevant information, respond with "No relevant information on this page."
                Include page number references in your response.
                """
                
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": "You are analyzing a document page by page."},
                        {"role": "user", "content": page_prompt}
                    ],
                    temperature=0,
                    max_tokens=500
                )
                
                page_content = response.choices[0].message.content
                tokens_used = response.usage.total_tokens
                add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                if "No relevant information on this page" not in page_content:
                    page_responses.append({
                        "page_number": page_num,
                        "content": page_content
                    })
                
                if st.session_state.debug_mode:
                    with st.expander(f"Debug: Page {page_num} Analysis", expanded=False):
                        st.write(f"Page {page_num} Content Length: {len(page['text'])}")
                        st.write("Relevant Content Found:" if page_responses else "No Relevant Content")
                        if page_responses:
                            st.write(page_content)
            
            if not page_responses:
                return "No relevant information found in the document to answer this question."
            
            # Combine all relevant page responses
            combined_response = "\n\n".join([
                f"From page {resp['page_number']}:\n{resp['content']}"
                for resp in page_responses
            ])
            
            return combined_response

        # Process query against stored documents
        # Initialize query variable from session state
        query = st.session_state.get('current_query', None)

        if query and (pdf_texts or ppt_texts) and not sales_data:
            try:
                responses = []
                
                # Process PPTs
                for ppt_name, ppt_data in [(ppt["name"], ppt) for ppt in ppt_texts]:
                    if st.session_state.debug_mode:
                        st.write(f"Analyzing PPT: {ppt_name}")
                    
                    try:
                        # Get document from database
                        doc_data = db.get_document(ppt_name)
                        if doc_data:
                            # Use slides/pages for PPT files
                            slides_data = {
                                "pages": doc_data["content"]["slides/pages"]  # Map slides to pages format
                            }
                            response = process_pdf_query_page_wise(query, slides_data)
                            
                            if "No relevant information found" not in response:
                                responses.append({
                                    "source": ppt_name,
                                    "type": "ppt",
                                    "content": response
                                })
                    except Exception as e:
                        if st.session_state.debug_mode:
                            st.error(f"Error processing PPT {ppt_name}: {str(e)}")
                        continue
                
                # Process PDFs
                for pdf_name, pdf_data in [(pdf["name"], pdf) for pdf in pdf_texts]:
                    if st.session_state.debug_mode:
                        st.write(f"Analyzing PDF: {pdf_name}")
                    
                    try:
                        # Get document from database
                        doc_data = db.get_document(pdf_name)
                        if doc_data:
                            response = process_pdf_query_page_wise(query, doc_data["content"])
                            
                            if response != "No relevant information found in the document to answer this question.":
                                responses.append({
                                    "pdf_name": pdf_name,
                                    "response": response
                                })
                    except Exception as e:
                        if st.session_state.debug_mode:
                            st.error(f"Error processing PDF {pdf_name}: {str(e)}")
                        continue

                # Synthesize all responses
                if responses:
                    try:
                        synthesis_prompt = f"""
                        Based on the following information from various sources:
                        {json.dumps(responses, indent=2)}
                        
                        Provide a comprehensive answer to: {query}
                        
                        Requirements:
                        1. Synthesize all relevant information
                        2. For each piece of information, cite its source (PPT or PDF filename)
                        3. Present information in a clear, logical flow
                        4. Focus only on relevant insights
                        5. Format the response as follows:
                            - Main answer with insights
                            - Sources section listing which files contributed to the answer
                        """
                        
                        final_response = client.chat.completions.create(
                            model="gpt-4o",
                            messages=[
                                {"role": "system", "content": "You are synthesizing information from multiple sources, providing clear answers with proper source citations."},
                                {"role": "user", "content": synthesis_prompt}
                            ],
                            temperature=0,
                            max_tokens=1000
                        )
                        tokens_used = final_response.usage.total_tokens
                        add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                        final_insight = final_response.choices[0].message.content
                        st.write(final_insight)
                        
                        # Update chat history
                        for chat in st.session_state['chat_history']:
                            if chat["query"] == query:
                                chat["response"] = final_insight
                                break
                        auto_save_chat_history()
                    except Exception as e:
                        if st.session_state.debug_mode:
                            st.error(f"Error synthesizing responses: {str(e)}")
                        st.error("An error occurred while processing your query. Please try again.")
                else:
                    st.write("No relevant information found in any of the uploaded documents.")

            except Exception as e:
                if st.session_state.debug_mode:
                    st.error(f"Error in query processing: {str(e)}")
                st.error("An error occurred while processing your query. Please try again.")
            finally:
                st.session_state.should_query = False
        elif query and sales_data:
            # Generate drill-down questions using the chat method
                try:
                        drill_questions = generate_drill_down_questions(temp_query, dataframes)
                        dril_results_PPT = []
                        if drill_questions['type'] == 'PPT':
                            st.write("### Drill Down Questions")
                            for dq in drill_questions['response']:
                                st.write(f"{dq}")
                            start_time = time.time()

                            for idx, dq in enumerate(drill_questions['response'], 1):
                                try:
                                    if 'state' in dq.lower():
                                        filtered_dataframes = [preprocess_data_for_state_query(df) for df in dataframes]
                                        temp_sales_data = SmartDatalake(filtered_dataframes, config={"llm": llm, "use_error_correction_framework": True, "data_dict": data_dict})
                                        dq_result = temp_sales_data.chat(dq)
                                    else:
                                        try:
                                            # st.write("Processing for ",dq)
                                            dq_result = sales_data.chat(dq)
                                        except e:
                                            st.error("Error: ",e)
                                        # st.write(f"Result for {idx} ",dq_result)
                                    if isinstance(dq_result,pd.DataFrame):
                                        dq_result = dq_result.to_string()
                                    else:
                                        dq_result = dq_result
                                    dril_results_PPT.append(dq_result)
                                
                                except Exception as e:
                                    st.error(f"Error processing drill-down question: {dq} -> {e}")
                            # st.write(drill_results)
                            total_time = time.time() - start_time
                            st.write(f"Total processing time for all questions: {total_time:.2f} seconds")

                            # Consolidate the drill-down results using the chat method
                            consolidation_prompt = consolidation_prompt.format(dril_results_PPT=dril_results_PPT,temp_query=temp_query)
                            response = client.chat.completions.create(
                                model="gpt-4o",
                                messages=[
                                    {"role": "system", "content": "You are an AI assistant."},
                                    {"role": "user", "content": consolidation_prompt}
                                ],
                                temperature=0,
                                max_tokens=1024
                            )
                            tokens_used = response.usage.total_tokens
                            add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                            final_insight = response.choices[0].message.content
                            st.write("### Final Consolidated Insight")
                            st.write(final_insight)

                            # Update chat history with the response
                            for chat in st.session_state['chat_history']:
                                if chat["query"] == temp_query:
                                    chat["response"] = final_insight
                                    break
                            auto_save_chat_history()
                            st.session_state.should_query = False
                        else:
                            st.write("### Drill Down Questions")
                            for dq in drill_questions['response']:
                                st.write(f"{dq}")

                            # Execute analysis for each drill-down question
                            drill_results = {
                                            "Overall Performance (2022-2024)": [],
                                            "Year-over-Year (YoY) Changes":[],
                                            "Regional Performance (2024)": [],
                                            "Market Penetration (2024)": [],
                                            "Consumer Analysis (2024)": [],
                                            # "Segment Performance (2024)": [],
                                            "Distribution (2024)": []
                                            }
                            keys = list(drill_results.keys())
                            start_time = time.time()
                            
                            for idx, dq in enumerate(drill_questions['response'], 1):
                                try:
                                    if 'state' in dq.lower():
                                        filtered_dataframes = [preprocess_data_for_state_query(df) for df in dataframes]
                                        temp_sales_data = SmartDatalake(filtered_dataframes, config={"llm": llm, "use_error_correction_framework": True, "data_dict": data_dict})
                                        dq_result = temp_sales_data.chat(dq)
                                    else:
                                        try:
                                            # st.write("Processing for ",dq)
                                            dq_result = sales_data.chat(dq)
                                        except e:
                                            st.error("Error: ",e)
                                        # st.write(f"Result for {idx}, {keys[idx-1]} ",dq_result)
                                    if isinstance(dq_result,pd.DataFrame):
                                        dq_result = dq_result.to_string()
                                    else:
                                        dq_result = dq_result
                                    drill_results[keys[idx-1]].append(dq_result)
                                
                                except Exception as e:
                                    st.error(f"Error processing drill-down question: {dq} -> {e}")
                            # st.write(drill_results)
                            total_time = time.time() - start_time
                            st.write(f"Total processing time for all questions: {total_time:.2f} seconds")

                            # Consolidate the drill-down results using the chat method
                            consolidation_prompt = consolidatin_prompt_main.format(temp_query=temp_query,drill_results=drill_results)
                            response = client.chat.completions.create(
                                model="gpt-4o",
                                messages=[
                                    {"role": "system", "content": "You are an AI assistant."},
                                    {"role": "user", "content": consolidation_prompt}
                                ],
                                temperature=0,
                                max_tokens=1024
                            )
                            final_insight = response.choices[0].message.content
                            tokens_used = response.usage.total_tokens
                            add_token_usage(tokens_used, st.session_state.last_query, usage_type="Unstructured Data")
                            st.write("### Final Consolidated Insight")
                            st.write(final_insight)

                            # Update chat history with the response
                            for chat in st.session_state['chat_history']:
                                if chat["query"] == temp_query:
                                    chat["response"] = final_insight
                                    break
                            auto_save_chat_history()
                            st.session_state.should_query = False
                except Exception as e:
                        st.error(f"An error occurred: {e}")
                        st.session_state.should_query = False
                finally:
                    st.session_state.should_query = False
        # ---- End your inline logic ----
        st.session_state.should_query_kb = False
# ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#This Part code is for the Triangulation.
    if st.session_state.triangulation:
        if not uploaded_files:
                st.warning("Please Upload files for Text extraction to process Triangulation analysis!")
        with st.spinner("Triangulation analysis..."):
            y_processor = st.session_state.query_processor
            result = y_processor.process_query(st.session_state.last_query)
            y_intent = result['intent']
            print(y_intent)
            add_token_usage(result['total_tokens'], st.session_state.last_query, usage_type="Intent Classification")

            engine = st.session_state.engine

            if y_intent == 'compare':
                init_sql, tokens_used_for_init_sql = engine.generate_compare_sql(result["original_query"])
            else:
                init_sql, tokens_used_for_init_sql = engine.generate_simple_sql(result["original_query"])
            # st.code(init_sql, language="sql")
            add_token_usage(tokens_used_for_init_sql, result["original_query"], usage_type="Initial SQL generation")

            # Validate/fix if desired
            validation_results = engine.validate_queries_with_explain([{"prompt": result["original_query"], "sql": init_sql}])
            print(validation_results)
            if validation_results["failed"]:
                fix_results = engine.fix_failed_queries_with_llm(validation_results["failed"])
                # print(fix_results)
                add_token_usage(fix_results['fixed_results'][0]['token_used'], result["original_query"], usage_type="For Fixing failed queries.")
                print(fix_results['fixed_results'][0]['llm_suggestion'])
                fixed_sql = fix_results['fixed_results'][0]['llm_suggestion']
                cleaned_query = fixed_sql.replace("```sql", "").replace("```", "").replace("\n", "").strip()
                try:
                    result_df = engine.execute_sql(cleaned_query)
                    st.session_state['last_generated_sql'] = cleaned_query
                    st.session_state['query_sql_map'][result["original_query"]] = cleaned_query
                    # st.dataframe(result_df)
                    # For Showing Initial Sql table
                    st.session_state['main_answer_df'] = result_df
                except Exception as e:
                    # st.error(f"Failed to execute SQL: {e}")
                    st.write("There's no data related to the question in the database.")
                
            elif validation_results["successful"]:
                try:
                    result_df = engine.execute_sql(init_sql)
                    st.session_state['last_generated_sql'] = init_sql
                    st.session_state['query_sql_map'][result["original_query"]] = init_sql
                    # st.dataframe(result_df)
                    # For Showing Initial Sql table
                    st.session_state['main_answer_df'] = result_df
                except Exception as e:
                    st.write("There's no data related to the question in the database.")

            if y_intent == 'compare' or y_intent == 'analysis':
                st.write('Text Extraction form the PDF')
                EXCEL_FILE_PATH = "callback_data.xlsx"
                user_defined_path = os.getcwd()
                udp = os.path.join(user_defined_path, "exports", "charts")

                if uploaded_files:
                    dataframes = []
                    pdf_texts = []
                    ppt_texts = []  # New list for PPT data
                    extracted_images = []
                    
                    for uploaded_file in uploaded_files:
                        try:
                            if uploaded_file.name.endswith('.pptx'):
                                # Process PPTX files
                                uploaded_file.seek(0)
                                pptx_result = test_process_pptx_with_gpt4v(uploaded_file)
                                add_token_usage(token_count=pptx_result['total_tokens'], query_text="Processing PPTX", usage_type="Processing PPTX")
                                if pptx_result:
                                    ppt_texts.append({
                                        "name": uploaded_file.name,
                                        "text": pptx_result["text"],
                                        "slides": pptx_result["slides/pages"]
                                    })
                                
                                # Extract images from PPTX
                                uploaded_file.seek(0)
                                pptx_images = test_extract_images_from_pptx(uploaded_file)
                                extracted_images.extend([{
                                    'source': uploaded_file.name,
                                    'type': 'pptx',
                                    **img
                                } for img in pptx_images])
                            
                            elif uploaded_file.name.endswith('.pdf'):
                                # Calculate file hash
                                file_content = uploaded_file.read()
                                file_hash = test_calculate_file_hash(file_content)
                                
                                # Check if file exists in database with same hash
                                doc_data = db.get_document(uploaded_file.name)
                                if doc_data and doc_data["file_hash"] == file_hash:
                                    if st.session_state.debug_mode:
                                        st.write(f"Using cached data from database for PDF: {uploaded_file.name}")
                                    pdf_texts.append({
                                        "name": uploaded_file.name,
                                        "text": doc_data["content"]["text"],
                                        "pages": doc_data["content"]["pages"]
                                    })
                                else:
                                    # Process new PDF file
                                    uploaded_file.seek(0)
                                    # Initial way to extracting data
                                    pdf_reader = PdfReader(uploaded_file)
                                    text = ""
                                    pages = []
                                    for i, page in enumerate(pdf_reader.pages):
                                        page_text = page.extract_text() or ""
                                        text += page_text
                                        pages.append({
                                            "page_number": i + 1,
                                            "text": page_text,
                                            "char_count": len(page_text),
                                            "word_count": len(page_text.split())
                                        })
                                    
                                    result = {
                                        "text": text,
                                        "pages": pages,
                                        "processed_date": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                                    }
                                    
                                    # Save to database
                                    db.save_document(
                                        filename=uploaded_file.name,
                                        file_type="pdf",
                                        content=result,
                                        file_hash=file_hash
                                    )
                                    
                                    pdf_texts.append({
                                        "name": uploaded_file.name,
                                        "text": text,
                                        "pages": pages
                                    })

                                # ------- Updated - 9/9/25 -------
                                uploaded_file.seek(0)  # Reset file pointer for image extraction

                                pdf_images = test_extract_images_from_pdf(uploaded_file)
                                extracted_images.extend([{
                                    'source': uploaded_file.name,
                                    'type': 'pdf',
                                    **img
                                } for img in pdf_images])

                                # ------- Till here -------

                            elif uploaded_file.name.endswith('.csv'):
                                df = pd.read_csv(uploaded_file)
                                dataframes.append(df)
                            elif uploaded_file.name.endswith('.xlsx'):
                                df = pd.read_excel(uploaded_file)
                                for col in df.select_dtypes(include=['object']).columns:
                                    unique_values_dict[col] = df[col].dropna().unique().tolist()
                                dataframes.append(df)
                            else:
                                st.error(f"Unsupported file format: {uploaded_file.name}")
                                continue
                            
                        except Exception as e:
                            st.error(f"Error loading file '{uploaded_file.name}': {e}")
                            continue

                    # Store extracted images in session state
                    st.session_state['extracted_images'] = extracted_images

                    # Process PDFs into chunks
                    if pdf_texts:
                        combined_text = "\n".join([pdf["text"] for pdf in pdf_texts])
                        text_splitter = RecursiveCharacterTextSplitter(
                            chunk_size=1000,
                            chunk_overlap=200,
                            length_function=len
                        )
                        chunks = text_splitter.split_text(combined_text)
                        st.session_state['pdf_chunks'] = chunks

                    if dataframes:
                        # st.write('dataframes')
                        sales_data = SmartDatalake(dataframes, config={"llm": llm, "use_error_correction_framework": True, "data_dict": data_dict})
                    else:
                        sales_data = None
                else:
                    sales_data = None

                query = st.session_state.get('current_query', None)

                if query and (pdf_texts or ppt_texts) and not sales_data:
                    try:
                        responses = []
                        
                        # Process PPTs
                        for ppt_name, ppt_data in [(ppt["name"], ppt) for ppt in ppt_texts]:
                            if st.session_state.debug_mode:
                                st.write(f"Analyzing PPT: {ppt_name}")
                            
                            try:
                                # Get document from database
                                doc_data = db.get_document(ppt_name)
                                if doc_data:
                                    # Use slides/pages for PPT files
                                    slides_data = {
                                        "pages": doc_data["content"]["slides/pages"]  # Map slides to pages format
                                    }
                                    response, total_tok = test_process_pdf_query_page_wise(query, slides_data)
                                    add_token_usage(token_count=total_tok, query_text="For processing pdf page wise in PPTX",usage_type="PDF page wise query in PPTX.")
                                    if "No relevant information found" not in response:
                                        responses.append({
                                            "source": ppt_name,
                                            "type": "ppt",
                                            "content": response
                                        })
                            except Exception as e:
                                if st.session_state.debug_mode:
                                    st.error(f"Error processing PPT {ppt_name}: {str(e)}")
                                continue
                        
                        # Process PDFs
                        for pdf_name, pdf_data in [(pdf["name"], pdf) for pdf in pdf_texts]:
                            if st.session_state.debug_mode:
                                st.write(f"Analyzing PDF: {pdf_name}")
                            
                            try:
                                # Get document from database
                                doc_data = db.get_document(pdf_name)
                                if doc_data:
                                    response, total_toke = test_process_pdf_query_page_wise(query, doc_data["content"])
                                    add_token_usage(token_count=total_toke, query_text="For processing pdf page wise",usage_type="PDF page wise query.")
                                    if response != "No relevant information found in the document to answer this question.":
                                        responses.append({
                                            "pdf_name": pdf_name,
                                            "response": response
                                        })
                            except Exception as e:
                                if st.session_state.debug_mode:
                                    st.error(f"Error processing PDF {pdf_name}: {str(e)}")
                                continue

                        # Synthesize all responses
                        if responses:
                            try:
                                synthesis_prompt = f"""
                                Based on the following information from various sources:
                                {json.dumps(responses, indent=2)}
                                
                                Provide a comprehensive answer to: {query}
                                
                                Requirements:
                                1. Synthesize all relevant information
                                2. For each piece of information, cite its source (PPT or PDF filename)
                                3. Present information in a clear, logical flow
                                4. Focus only on relevant insights
                                5. Format the response as follows:
                                    - Main answer with insights
                                    - Sources section listing which files contributed to the answer
                                """
                                
                                final_response = client.chat.completions.create(
                                    model="gpt-4o",
                                    messages=[
                                        {"role": "system", "content": "You are synthesizing information from multiple sources, providing clear answers with proper source citations."},
                                        {"role": "user", "content": synthesis_prompt}
                                    ],
                                    temperature=0,
                                    max_tokens=1000
                                )
                                tokens_used = final_response.usage.total_tokens
                                add_token_usage(tokens_used, "final response", usage_type="Unstructured Data")
                                final_insight = final_response.choices[0].message.content
                                st.write(final_insight)
                                print(final_response.usage, 'Final response')
                                # # Update chat history
                                # for chat in st.session_state['chat_history']:
                                #     if chat["query"] == query:
                                #         chat["response"] = final_insight
                                #         break
                                # auto_save_chat_history()
                            except Exception as e:
                                if st.session_state.debug_mode:
                                    st.error(f"Error synthesizing responses: {str(e)}")
                                st.error("An error occurred while processing your query. Please try again.")
                        else:
                            st.write("No relevant information found in any of the uploaded documents.")

                    except Exception as e:
                        if st.session_state.debug_mode:
                            st.error(f"Error in query processing: {str(e)}")
                        st.error("An error occurred while processing your query. Please try again.")
                    finally:
                        st.session_state.should_query = False
        st.session_state.triangulation = None
# Till here
# ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    if st.session_state.should_run_filtered:
        with st.spinner("Processing filtered analysis..."):
            # Before filtering the result, I want give sql table for the master query.
            y_processor = st.session_state.query_processor
            result = y_processor.process_query(st.session_state.last_query)
            y_intent = result['intent']
            add_token_usage(result['total_tokens'], st.session_state.last_query, usage_type="Intent Classification")

            engine = st.session_state.engine

            if y_intent == 'compare':
                init_sql, tokens_used_for_init_sql = engine.generate_compare_sql(result["original_query"])
            else:
                init_sql, tokens_used_for_init_sql = engine.generate_simple_sql(result["original_query"])
            # st.code(init_sql, language="sql")
            add_token_usage(tokens_used_for_init_sql, result["original_query"], usage_type="Initial SQL generation")

            # Validate/fix if desired
            validation_results = engine.validate_queries_with_explain([{"prompt": result["original_query"], "sql": init_sql}])
            print(validation_results)
            if validation_results["failed"]:
                fix_results = engine.fix_failed_queries_with_llm(validation_results["failed"])
                # print(fix_results)
                add_token_usage(fix_results['fixed_results'][0]['token_used'], result["original_query"], usage_type="For Fixing failed queries.")
                print(fix_results['fixed_results'][0]['llm_suggestion'])
                fixed_sql = fix_results['fixed_results'][0]['llm_suggestion']
                cleaned_query = fixed_sql.replace("```sql", "").replace("```", "").replace("\n", "").strip()
                try:
                    result_df = engine.execute_sql(cleaned_query)
                    st.session_state['last_generated_sql'] = cleaned_query
                    st.session_state['query_sql_map'][result["original_query"]] = cleaned_query
                    # st.dataframe(result_df)
                    # For Showing Initial Sql table
                    st.session_state['main_answer_df'] = result_df
                except Exception as e:
                    # st.error(f"Failed to execute SQL: {e}")
                    st.write("There's no data related to the question in the database.")
                
            elif validation_results["successful"]:
                try:
                    result_df = engine.execute_sql(init_sql)
                    st.session_state['last_generated_sql'] = init_sql
                    st.session_state['query_sql_map'][result["original_query"]] = init_sql
                    # st.dataframe(result_df)
                    # For Showing Initial Sql table
                    st.session_state['main_answer_df'] = result_df
                except Exception as e:
                    st.write("There's no data related to the question in the database.")

            # -------------------------------------------------------------------------------
            prompts_to_drill = st.session_state.database_previous_prompts
            filtered_result = engine.handle_prompt(st.session_state.last_query, intent=y_intent, previous_prompts=prompts_to_drill[:-1])
            print("###############################################")
            print(filtered_result)
            if y_intent != 'simple' and y_intent != 'compare':
                add_token_usage(filtered_result['total_tokens'], st.session_state.last_query, usage_type="Drill down question with SQL generation")

            print("###############################################")
            # Store results in session state for next step
            if y_intent == 'simple':
                st.session_state.filtered_results = filtered_result
                st.session_state.intent_type = "simple"
                st.session_state.should_run_filtered = False
            if y_intent == 'compare':
                st.session_state.filtered_results = filtered_result
                st.session_state.intent_type = "compare"
                st.session_state.should_run_filtered = False
            elif y_intent == 'analysis':
                if hasattr(filtered_result, 'get') and "sub_prompts_sql" in filtered_result:
                    validation = engine.validate_queries_with_explain(filtered_result["sub_prompts_sql"])
                else:
                    validation = {"successful": [], "failed": []}
                st.session_state.filtered_results = validation
                st.session_state.intent_type = "analysis"
                st.session_state.should_run_filtered = False
                st.session_state.should_show_selection = True

    # --- Always Show Main Answer Table --- #
    if st.session_state['main_answer_df'] is not None:
        
        df_to_show = st.session_state['main_answer_df'].copy()
        # Format numeric columns with commas and two decimals
        df = df_to_show.copy()

        df = df.reset_index(drop=True)

        def format_number(x, is_year=False):
            if pd.isna(x):
                return ""
            
            if is_year:
                return f"{int(x)}"
            
            # If it's effectively an integer (e.g., 2023.0), show without decimals
            if isinstance(x, (int, float)):
                if float(x).is_integer():
                    return f"{int(x)}"
                else:
                    return f"{x:,.2f}"
            
            return x  # For non-numeric types

        # Apply formatting
        for col in df.columns:
            if col.lower() == "year":
                df[col] = df[col].map(lambda x: format_number(x, is_year=True))
            elif pd.api.types.is_numeric_dtype(df[col]):
                df[col] = df[col].map(format_number)

        # Show without index
        st.dataframe(df, hide_index=True)


    # ----- For Analysis Intent: Show Checklist ----
    if st.session_state.get("should_show_selection") and st.session_state.get("filtered_results"):
        engine = st.session_state.engine
        successful = st.session_state.filtered_results.get("successful", []).copy() # ensure copy, avoid mutation
        failed = st.session_state.filtered_results.get("failed", [])
        newly_fixed = []
        if failed:
            print('About to fix failed queries. Failed=', failed)
            print('Type:', type(failed))
            fix_results = engine.fix_failed_queries_with_llm(failed)
            print(fix_results, 'Fixing Failed Queries %^&*@#')
            for fix in fix_results["fixed_results"]:
                if fix.get('explain_success'):
                    fixed_prompt = {
                        "prompt": fix["prompt"],
                        "sql": fix["llm_suggestion"]
                    }
                    successful.append(fixed_prompt)
                    newly_fixed.append(fixed_prompt)
                    add_token_usage(fix.get('token_used', 0), fix["prompt"], usage_type="For Fixing failed queries.")
                else:
                    st.warning(f"Could not automatically fix failed query: {fix['prompt']}. Please review manually.")

        # -- Store updated list for the next steps --
        prompts = [item["prompt"] for item in successful]
        options = prompts + ["Query Knowledge Base"]
        st.session_state.all_successful_for_selection = successful  # Save for access in next UI section

        st.markdown("### Select one or more analyses to run:")

        # Keep checked state in session
        if "checkbox_states" not in st.session_state or \
                st.session_state.get("checkbox_option_set") != set(options):
            # New question or changed set: reset the checkboxes
            st.session_state.checkbox_states = [False] * len(options)
            st.session_state.checkbox_option_set = set(options)

        # Render checkboxes with persistent state
        checked = []
        for i, option in enumerate(options):
            st.session_state.checkbox_states[i] = st.checkbox(
                option,
                value=st.session_state.checkbox_states[i],
                key=f"chk_{i}"
            )
            if st.session_state.checkbox_states[i]:
                checked.append(option)

        submit_btn = st.button("Submit", key="submit_filtered_btn", disabled=len(checked) == 0)

        if submit_btn:
            st.session_state.selected_prompts = list(checked)
            st.session_state.should_show_selection = False  # hide checklist after submit

    # ----- Show Final Filtered and/or Unfiltered Answers -----
    if st.session_state.get("selected_prompts") is not None:
        all_success = st.session_state.get('all_successful_for_selection') or \
        st.session_state.filtered_results.get("successful", []) or []
        prompt_to_sql = {item["prompt"]: item["sql"] for item in all_success}

        # result = y_processor.process_query(st.session_state.last_query)
        # y_intent = result['intent']

        engine = st.session_state.engine

        for prompt in st.session_state.selected_prompts:
            if prompt != "Query Knowledge Base":
                st.markdown(f"**Filtered analysis for:** {prompt}")
                sql = prompt_to_sql.get(prompt, "No SQL for this prompt.")
                # st.code(sql, language="sql")
                # Here you can run/display SQL result as needed
                # Validate/fix if desired
                # ================================
                # For Validation
                validation_results = engine.validate_queries_with_explain([{"prompt": prompt, "sql": sql}])
                print(validation_results)

                if validation_results["failed"]:
                    fix_results = engine.fix_failed_queries_with_llm(validation_results["failed"])
                    # print(fix_results)
                    add_token_usage(fix_results['fixed_results'][0]['token_used'], prompt, usage_type="For Fixing failed queries.")
                    sql_cleaned_query = sql.replace("```sql", "").replace("```", "").replace("\n", "").strip()
                    result_df = engine.execute_sql(sql_cleaned_query)
                    st.session_state['query_sql_map'][prompt] = sql_cleaned_query
                    df = result_df.copy()
                    # Reset index to remove index column from output
                    df = df.reset_index(drop=True)

                    def format_number(x, is_year=False):
                        if pd.isna(x):
                            return ""
                        if is_year:
                            return f"{int(x)}"
                        else:
                            return f"{x:,.2f}"

                    for col in df.columns:
                        if col.lower() == "year":
                            df[col] = df[col].map(lambda x: format_number(x, is_year=True))
                        elif pd.api.types.is_numeric_dtype(df[col]):
                            df[col] = df[col].map(format_number)

                    st.dataframe(df, hide_index=True)
                    
                elif validation_results["successful"]:
                    result_df = engine.execute_sql(sql)
                    st.session_state['query_sql_map'][prompt] = sql
                    df = result_df.copy()

                    # Reset index to remove index column from output
                    df = df.reset_index(drop=True)

                    def format_number(x, is_year=False):
                        if pd.isna(x):
                            return ""
                        if is_year:
                            return f"{int(x)}"
                        else:
                            return f"{x:,.2f}"

                    for col in df.columns:
                        if col.lower() == "year":
                            df[col] = df[col].map(lambda x: format_number(x, is_year=True))
                        elif pd.api.types.is_numeric_dtype(df[col]):
                            df[col] = df[col].map(format_number)

                    st.dataframe(df, hide_index=True)
                    # sql_cleaned_query = sql.replace("```sql", "").replace("```", "").replace("\n", "").strip()
                    # result_df = engine.execute_sql(sql_cleaned_query)
                    # st.dataframe(result_df)
                # ===============================
        if "Query Knowledge Base" in st.session_state.selected_prompts:
            st.markdown("---\n**Query Knowledge Base:**")
            # Add your normal app.py fallback answer logic here
            EXCEL_FILE_PATH = "callback_data.xlsx"
            user_defined_path = os.getcwd()
            udp = os.path.join(user_defined_path, "exports", "charts")

            # Function to parse callback output
            def parse_callback_output(cb) -> dict:
                output_lines = str(cb).split("\n")
                parsed_output = {
                    "total_tokens": int(output_lines[0].split(":")[1].strip()),
                    "prompt_tokens": int(output_lines[1].split(":")[1].strip()),
                    "completion_tokens": int(output_lines[2].split(":")[1].strip()),
                    "total_cost": float(output_lines[3].split(":")[1].strip().replace("$", "").strip()),
                }
                return parsed_output

            # Function to append data to Excel
            def append_to_excel(file_path: str, data: dict):
                if not os.path.exists(file_path):
                    wb = Workbook()
                    ws = wb.active
                    ws.append(["Total Tokens", "Prompt Tokens", "Completion Tokens", "Total Cost (USD)"])
                    wb.save(file_path)
                wb = openpyxl.load_workbook(file_path)
                ws = wb.active
                ws.append([data["total_tokens"], data["prompt_tokens"], data["completion_tokens"], data["total_cost"]])
                wb.save(file_path)


            # Add these new functions after the imports
            def extract_images_from_pdf(pdf_file):
                """Extract images and charts from PDF using PyMuPDF with focus on chart detection"""
                images = []
                temp_file = None
                doc = None
                
                try:
                    # Create a temporary file to save the uploaded PDF
                    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                    temp_file_path = temp_file.name
                    temp_file.write(pdf_file.read())
                    temp_file.close()
                    
                    # Open the PDF file
                    doc = fitz.open(temp_file_path)
                    
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        
                        # Method 1: Standard image extraction
                        image_list = page.get_images()
                        for img_index, img in enumerate(image_list):
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            # Convert to PIL Image
                            image = Image.open(io.BytesIO(image_bytes))
                            if is_valid_image_size(image):
                                images.append({
                                    'page': page_num + 1,
                                    'index': img_index,
                                    'image': image,
                                    'format': base_image["ext"],
                                    'method': 'direct_extraction',
                                    'width': image.width,
                                    'height': image.height
                                })
                        
                        # Method 2: Extract potential chart regions by analyzing page content
                        # Get all drawings on the page
                        paths = page.get_drawings()
                        if paths:
                            # Extract regions with significant drawing content
                            regions = analyze_drawing_regions(paths)
                            for region_index, region in enumerate(regions):
                                try:
                                    # Render the region with higher resolution
                                    zoom = 2  # Increase if needed for better quality
                                    mat = fitz.Matrix(zoom, zoom)
                                    pix = page.get_pixmap(matrix=mat, clip=region)
                                    
                                    # Convert to PIL Image
                                    img_data = Image.frombytes("RGB", [pix.width, pix.height], pix.tobytes())
                                    
                                    if is_valid_image_size(img_data):
                                        images.append({
                                            'page': page_num + 1,
                                            'index': f'chart_{region_index}',
                                            'image': img_data,
                                            'format': 'png',
                                            'method': 'region_extraction',
                                            'width': img_data.width,
                                            'height': img_data.height,
                                            'bbox': [region.x0, region.y0, region.x1, region.y1]
                                        })
                                except Exception as e:
                                    print(f"Error extracting region: {e}")
                    
                    return images
                    
                except Exception as e:
                    st.error(f"Error extracting images from PDF: {e}")
                    return []
                
                finally:
                    if doc:
                        doc.close()
                    if temp_file and os.path.exists(temp_file.name):
                        try:
                            os.unlink(temp_file.name)
                        except Exception as e:
                            print(f"Error removing temporary file: {e}")

            def is_valid_image_size(image):
                """Check if image meets minimum size requirements"""
                MIN_WIDTH = 100
                MIN_HEIGHT = 100
                MAX_DIMENSION = 4000  # Prevent extremely large images
                
                return (MIN_WIDTH <= image.width <= MAX_DIMENSION and 
                        MIN_HEIGHT <= image.height <= MAX_DIMENSION)

            def analyze_drawing_regions(paths):
                """Analyze drawing paths to identify potential chart regions"""
                regions = []
                
                if not paths:
                    return regions
                
                # Group nearby paths that might form a chart
                current_group = []
                
                for path in paths:
                    rect = path['rect']  # Get the bounding rectangle
                    
                    # Skip tiny drawings
                    if rect.width < 20 or rect.height < 20:
                        continue
                    
                    if current_group:
                        # Check if this path is close to the current group
                        last_rect = current_group[-1]['rect']
                        if (abs(rect.x0 - last_rect.x1) < 50 and 
                            abs(rect.y0 - last_rect.y1) < 50):
                            current_group.append(path)
                else:
                        current_group.append(path)
                
                # Don't forget the last group
                if len(current_group) >= 3:
                    merged_rect = merge_rects([p['rect'] for p in current_group])
                    if merged_rect.width >= 100 and merged_rect.height >= 100:
                        regions.append(merged_rect)
                
                return regions

            def merge_rects(rects):
                """Merge multiple rectangles into one bounding rectangle"""
                if not rects:
                    return None
                
                x0 = min(rect.x0 for rect in rects)
                y0 = min(rect.y0 for rect in rects)
                x1 = max(rect.x1 for rect in rects)
                y1 = max(rect.y1 for rect in rects)
                
                # Add some padding
                padding = 10
                return fitz.Rect(x0 - padding, y0 - padding, x1 + padding, y1 + padding)

            def extract_images_from_pptx(pptx_file):
                """Extract images from PowerPoint file"""
                images = []
                try:
                    prs = Presentation(pptx_file)
                    
                    for slide_num, slide in enumerate(prs.slides):
                        for shape_num, shape in enumerate(slide.shapes):
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                image_stream = io.BytesIO(shape.image.blob)
                                image = Image.open(image_stream)
                                images.append({
                                    'slide': slide_num + 1,
                                    'index': shape_num,
                                    'image': image,
                                    'format': image.format.lower()
                                })
                    return images
                except Exception as e:
                    st.error(f"Error extracting images from PPTX: {e}")
                    return []

            def convert_pptx_to_pdf(input_dir, output_dir):
                """Convert PPTX to PDF using platform-specific methods"""
                try:
                    if os.name == 'nt':  # Windows
                        # Import pythoncom only for Windows
                        import pythoncom
                        # Initialize COM for this thread
                        pythoncom.CoInitialize()
                        try:
                            from pptxtopdf import convert
                            convert(input_dir, output_dir)
                        finally:
                            # Uninitialize COM
                            pythoncom.CoUninitialize()
                    else:  # Linux/Unix
                        # Check if libreoffice is installed
                        try:
                            import subprocess
                            subprocess.run(['libreoffice', '--version'], capture_output=True, check=True)
                        except (subprocess.CalledProcessError, FileNotFoundError):
                            raise RuntimeError("LibreOffice is not installed. Please install it using: sudo apt-get install libreoffice")
                        
                        # Get the input file path
                        input_file = os.path.join(input_dir, os.listdir(input_dir)[0])  # Get the first file in the directory
                        
                        # Convert using libreoffice with explicit file path
                        subprocess.run([
                            'libreoffice',
                            '--headless',
                            '--convert-to', 'pdf',
                            '--outdir', output_dir,
                            input_file
                        ], check=True)
                        
                        # Add debug logging
                        if st.session_state.debug_mode:
                            st.write(f"Input file: {input_file}")
                            st.write(f"Output directory: {output_dir}")
                            st.write(f"Files in output directory: {os.listdir(output_dir)}")
                        
                except Exception as e:
                    raise RuntimeError(f"Error converting PPTX to PDF: {str(e)}")

            def process_pptx_with_gpt4v(pptx_file):
                """Process PPTX using GPT-4 Vision by first converting to PDF then processing like a PDF"""
                # Calculate file hash
                file_content = pptx_file.read()
                file_hash = calculate_file_hash(file_content)
                
                # Check if file exists in database with same hash
                doc_data = db.get_document(pptx_file.name)
                if doc_data and doc_data["file_hash"] == file_hash:
                    if st.session_state.debug_mode:
                        st.write(f"Using cached data from database for PPT: {pptx_file.name}")
                    return doc_data["content"]
                
                # Reset file pointer for processing
                pptx_file.seek(0)
                
                # Create temporary directories
                temp_input_dir = tempfile.mkdtemp()
                temp_output_dir = tempfile.mkdtemp()
                
                try:
                    # Save uploaded file to temp directory
                    temp_input_path = os.path.join(temp_input_dir, pptx_file.name)
                    with open(temp_input_path, 'wb') as f:
                        f.write(file_content)
                    
                    # Convert PPTX to PDF using platform-specific method
                    convert_pptx_to_pdf(temp_input_dir, temp_output_dir)
                    
                    # Find the converted PDF file
                    pdf_filename = os.path.splitext(pptx_file.name)[0] + '.pdf'
                    pdf_path = os.path.join(temp_output_dir, pdf_filename)
                    
                    if not os.path.exists(pdf_path):
                        raise FileNotFoundError(f"PDF file not found at expected location: {pdf_path}")
                    
                    if st.session_state.debug_mode:
                        st.write(f"Successfully converted PPTX to PDF: {pdf_path}")
                    
                    # Process the PDF using our existing PDF processing logic
                    with open(pdf_path, 'rb') as pdf_file:
                        # Initialize OpenAI client
                        # client = OpenAIClient()
                        
                        # Convert PDF pages to images
                        images = convert_from_path(pdf_path)
                        page_texts = []
                        
                        for i, image in enumerate(images):
                            # Convert image to base64
                            buffered = BytesIO()
                            image.save(buffered, format="PNG")
                            base64_image = base64.b64encode(buffered.getvalue()).decode()
                            
                            # Extract text using GPT-4V
                            response = client.chat.completions.create(
                                model="gpt-4o",
                                messages=[
                                    {
                                        "role": "user",
                                        "content": [
                                            {
                                                "type": "text",
                                                "text": "Extract all content from this slide. Include text, describe any images, charts, or diagrams, and maintain the formatting and structure."
                                            },
                                            {
                                                "type": "image_url",
                                                "image_url": {
                                                    "url": f"data:image/png;base64,{base64_image}",
                                                    "detail": "high"
                                                }
                                            }
                                        ]
                                    }
                                ],
                                max_tokens=1000
                            )
                            
                            extracted_text = response.choices[0].message.content
                            page_texts.append({
                                "slide_number": i + 1,
                                "text": extracted_text,
                                "char_count": len(extracted_text),
                                "word_count": len(extracted_text.split())
                            })
                            
                            if st.session_state.debug_mode:
                                with st.expander(f"Debug: Extracted Content Slide {i+1}", expanded=False):
                                    st.write("#### Extracted Content")
                                    st.text_area("Text Content", extracted_text, height=300)
                                    st.write(f"Characters extracted: {len(extracted_text)}")
                                    st.write(f"Words extracted: {len(extracted_text.split())}")
                        
                        # Combine all slide texts
                        combined_text = "\n\n".join([page["text"] for page in page_texts])
                        result = {
                            "text": combined_text,
                            "slides/pages": page_texts,
                            "processed_date": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                        }
                        
                        # Save to database
                        db.save_document(
                            filename=pptx_file.name,
                            file_type="pptx",
                            content=result,
                            file_hash=file_hash
                        )
                        
                        return result
                        
                except Exception as e:
                    if st.session_state.debug_mode:
                        st.error(f"Error during PDF conversion/processing: {str(e)}")
                    raise e
                    
                finally:
                    # Clean up temporary directories
                    try:
                        if os.path.exists(temp_input_dir):
                            shutil.rmtree(temp_input_dir)
                        if os.path.exists(temp_output_dir):
                            shutil.rmtree(temp_output_dir)
                    except Exception as e:
                        if st.session_state.debug_mode:
                            st.error(f"Error cleaning up temporary files: {str(e)}")

            # Add new function for GPT-4 Vision extraction
            def extract_text_with_gpt4v(pdf_file):
                """Extract text from PDF using GPT-4 Vision"""
                try:
                    # Create a temporary file to save the uploaded PDF
                    temp_file = None
                    text_results = []
                    
                    try:
                        # Save uploaded file to temporary file
                        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                        temp_file.write(pdf_file.read())
                        temp_file.close()
                        
                        try:
                            # Convert PDF pages to images using pdf2image
                            images = convert_from_path(temp_file.name)
                            
                            # Process each page
                            for i, image in enumerate(images):
                                if st.session_state.debug_mode:
                                    with st.expander(f"Debug: GPT-4V Processing Page {i+1}", expanded=False):
                                        st.write("---")
                                        st.write(f"### Processing Page {i+1}")
                                        st.image(image, caption=f"Page {i+1}", use_column_width=True)
                                
                                # Convert PIL Image to bytes
                                img_byte_arr = io.BytesIO()
                                image.save(img_byte_arr, format='PNG')
                                img_byte_arr = img_byte_arr.getvalue()
                                
                                # Encode image to base64
                                base64_image = base64.b64encode(img_byte_arr).decode('utf-8')
                                
                                # Create GPT-4V prompt with correct image URL format
                                response = client.chat.completions.create(
                                    model="gpt-4o",
                                    messages=[
                                        {
                                            "role": "user",
                                            "content": [
                                                {
                                                    "type": "text",
                                                    "text": "Extract all text content from this image. Preserve formatting and structure. Include any relevant tables, charts, or diagrams descriptions."
                                                },
                                                {
                                                    "type": "image_url",
                                                    "image_url": {
                                                        "url": f"data:image/png;base64,{base64_image}",
                                                        "detail": "high"
                                                    }
                                                }
                                            ]
                                        }
                                    ],
                                    max_tokens=1000
                                )
                                
                                extracted_text = response.choices[0].message.content
                                
                                if st.session_state.debug_mode:
                                    with st.expander(f"Debug: Extracted Content Page {i+1}", expanded=False):
                                        st.write("#### Extracted Content")
                                        st.text_area("Text Content", extracted_text, height=300)
                                        st.write(f"Characters extracted: {len(extracted_text)}")
                                        st.write(f"Words extracted: {len(extracted_text.split())}")
                                
                                text_results.append(extracted_text)
                            
                            return "\n\n".join(text_results)
                            
                        except Exception as e:
                            if st.session_state.debug_mode:
                                st.error(f"Error processing PDF: {str(e)}")
                            raise e
                            
                    finally:
                        # Clean up temporary file
                        if temp_file and os.path.exists(temp_file.name):
                            try:
                                os.unlink(temp_file.name)
                            except Exception as e:
                                if st.session_state.debug_mode:
                                    st.error(f"Error removing temporary file: {e}")
                            
                except Exception as e:
                    st.error(f"Error in GPT-4V processing: {e}")
                    if st.session_state.debug_mode:
                        st.error(f"Detailed error information: {str(e)}")
                    return ""

            # Prepare datasets for SmartDatalake and process PDFs immediately
            if uploaded_files:
                dataframes = []
                pdf_texts = []
                ppt_texts = []  # New list for PPT data
                extracted_images = []
                
                for uploaded_file in uploaded_files:
                    try:
                        if uploaded_file.name.endswith('.pptx'):
                            # Process PPTX files
                            uploaded_file.seek(0)
                            pptx_result = process_pptx_with_gpt4v(uploaded_file)
                            if pptx_result:
                                ppt_texts.append({
                                    "name": uploaded_file.name,
                                    "text": pptx_result["text"],
                                    "slides": pptx_result["slides/pages"]
                                })
                            
                            # Extract images from PPTX
                            uploaded_file.seek(0)
                            pptx_images = extract_images_from_pptx(uploaded_file)
                            extracted_images.extend([{
                                'source': uploaded_file.name,
                                'type': 'pptx',
                                **img
                            } for img in pptx_images])
                        
                        elif uploaded_file.name.endswith('.pdf'):
                            # Calculate file hash
                            file_content = uploaded_file.read()
                            file_hash = calculate_file_hash(file_content)
                            
                            # Check if file exists in database with same hash
                            doc_data = db.get_document(uploaded_file.name)
                            if doc_data and doc_data["file_hash"] == file_hash:
                                if st.session_state.debug_mode:
                                    st.write(f"Using cached data from database for PDF: {uploaded_file.name}")
                                pdf_texts.append({
                                    "name": uploaded_file.name,
                                    "text": doc_data["content"]["text"],
                                    "pages": doc_data["content"]["pages"]
                                })
                            else:
                                # Process new PDF file
                                uploaded_file.seek(0)
                                pdf_reader = PdfReader(uploaded_file)
                                text = ""
                                pages = []
                                for i, page in enumerate(pdf_reader.pages):
                                    page_text = page.extract_text() or ""
                                    text += page_text
                                    pages.append({
                                        "page_number": i + 1,
                                        "text": page_text,
                                        "char_count": len(page_text),
                                        "word_count": len(page_text.split())
                                    })
                                
                                result = {
                                    "text": text,
                                    "pages": pages,
                                    "processed_date": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                                }
                                
                                # Save to database
                                db.save_document(
                                    filename=uploaded_file.name,
                                    file_type="pdf",
                                    content=result,
                                    file_hash=file_hash
                                )
                                
                                pdf_texts.append({
                                    "name": uploaded_file.name,
                                    "text": text,
                                    "pages": pages
                                })
                        
                        elif uploaded_file.name.endswith('.csv'):
                            df = pd.read_csv(uploaded_file)
                            dataframes.append(df)
                        elif uploaded_file.name.endswith('.xlsx'):
                            df = pd.read_excel(uploaded_file)
                            for col in df.select_dtypes(include=['object']).columns:
                                unique_values_dict[col] = df[col].dropna().unique().tolist()
                            dataframes.append(df)
                        else:
                            st.error(f"Unsupported file format: {uploaded_file.name}")
                            continue
                        
                    except Exception as e:
                        st.error(f"Error loading file '{uploaded_file.name}': {e}")
                        continue

                # Store extracted images in session state
                st.session_state['extracted_images'] = extracted_images

                # Process PDFs into chunks
                if pdf_texts:
                    combined_text = "\n".join([pdf["text"] for pdf in pdf_texts])
                    text_splitter = RecursiveCharacterTextSplitter(
                        chunk_size=1000,
                        chunk_overlap=200,
                        length_function=len
                    )
                    chunks = text_splitter.split_text(combined_text)
                    st.session_state['pdf_chunks'] = chunks

                if dataframes:
                    # st.write('dataframes')
                    sales_data = SmartDatalake(dataframes, config={"llm": llm, "use_error_correction_framework": True, "data_dict": data_dict})
                else:
                    sales_data = None
            else:
                sales_data = None


            # Function to generate up to 3 drill-down questions using the LLM's chat method
            def generate_drill_down_questions(query, dataframes):
                # Check if we have brand metrics data by looking at the dataframe structure
                def is_brand_metrics_df(df):
                    required_columns = ['Brand', 'Date', 'TOM', 'Spontaneous_Awareness', 'MOUB', 'Purchase_6M']
                    return all(col in df.columns for col in required_columns)
                
                # Determine the type of analysis based on the data structure
                has_brand_metrics = any(is_brand_metrics_df(df) for df in dataframes)
                
                if has_brand_metrics:
                    drill_prompt_ppt = dril_ppt.format(query=query)
                    response_PPT = client.chat.completions.create(
                            model="gpt-4o",
                            messages=[
                                {"role": "system", "content": "You are an AI assistant."},
                                {"role": "user", "content": drill_prompt_ppt}
                            ],
                            temperature=0,
                            max_tokens=512
                        )
                    drill_response_PPT = response_PPT.choices[0].message.content
                    clean_response_PPT = drill_response_PPT.replace("```python", "").replace("```", "").strip()
                    parsed_questions_PPT= ast.literal_eval(clean_response_PPT)
                    return {"type":'PPT',"response":parsed_questions_PPT}
                else:
                    
                    drill_prompt_main = drill_p_m.format(query=query,
                                                            data_dictt=json.dumps(data_dict, indent=2),
                                                            unique_values_dict=json.dumps({col: unique_values_dict[col] for col in unique_values_dict if col in data_dict}, indent=2))
                    response_main = client.chat.completions.create(
                            model="gpt-4o",
                            messages=[
                                {"role": "system", "content": "You are an AI assistant."},
                                {"role": "user", "content": drill_prompt_main}
                            ],
                            temperature=0,
                            max_tokens=512
                        )
                
                    drill_response_main = response_main.choices[0].message.content
                    clean_response_main = drill_response_main.replace("```python", "").replace("```", "").strip()
                    parsed_questions_main = ast.literal_eval(clean_response_main)
                    return {"type":'main',"response":parsed_questions_main}

            # Function to auto-save chat history
            def auto_save_chat_history():
                """Automatically save the chat history after each interaction"""
                if st.session_state['chat_history']:
                    # Create a 'chat_history' directory if it doesn't exist
                    os.makedirs('chat_history', exist_ok=True)
                    
                    # Use a fixed filename with current date
                    current_date = datetime.now().strftime("%Y%m%d")
                    filename = f"chat_history/chat_history_{current_date}.json"
                    
                    # Add timestamp to each chat entry that doesn't have one
                    chat_history_with_timestamp = []
                    for chat in st.session_state['chat_history']:
                        chat_entry = chat.copy()
                        if 'timestamp' not in chat_entry:
                            chat_entry['timestamp'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                        chat_history_with_timestamp.append(chat_entry)
                        
                    try:
                        with open(filename, 'w', encoding='utf-8') as f:
                            json.dump(chat_history_with_timestamp, f, indent=2, ensure_ascii=False)
                    except Exception as e:
                        st.error(f"Error auto-saving chat history: {e}")

            # Add this function to filter data before querying
            def preprocess_data_for_state_query(df):
                if 'Region' in df.columns:
                    # Keep all rows except those where MARKET starts with 'ALL INDIA'
                    return df[~df['Region'].str.startswith('ALL INDIA', na=False)]
                return df

            # Function to retrieve relevant chunks based on the query with error handling
            def retrieve_relevant_chunks(query, chunks, top_k=3):
                if not chunks:
                    st.error("No chunks available to process the query.")
                    return []
                prompt = f"""
                Given the query: "{query}", and the following document chunks:
                {json.dumps(chunks[:5], indent=2)} (showing first 5 chunks for brevity),
                identify the top {top_k} most relevant chunks to answer the query.
                Return the indices of these chunks as a list (e.g., [0, 2, 4]).
                """
                try:
                    response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": "You are an AI assistant that identifies relevant text chunks."},
                            {"role": "user", "content": prompt}
                        ],
                        temperature=0,
                        max_tokens=100
                    )
                    response_content = response.choices[0].message.content.strip()
                    indices = ast.literal_eval(response_content)
                    if not isinstance(indices, list):
                        raise ValueError("LLM did not return a valid list of indices")
                    return [chunks[i] for i in indices if i < len(chunks)]
                except Exception as e:
                    print(f"Error retrieving relevant chunks: {e}")
                    return chunks[:top_k]  # Fallback to first top_k chunks

            def process_pdf_query_page_wise(query, pdf_data):
                """Process query by analyzing each page's content individually"""
                page_responses = []
                
                # First, analyze each page individually
                for page in pdf_data['pages']:
                    # Handle both page_number and slide_number fields
                    page_num = page.get('page_number', page.get('slide_number', 1))
                    page_prompt = f"""
                    Analyzing page {page_num} of the document.
                    Content: {page['text']}
                    
                    Question: {query}
                    
                    If this page contains relevant information to answer the question, provide the relevant details.
                    If this page does not contain relevant information, respond with "No relevant information on this page."
                    Include page number references in your response.
                    """
                    
                    response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": "You are analyzing a document page by page."},
                            {"role": "user", "content": page_prompt}
                        ],
                        temperature=0,
                        max_tokens=500
                    )
                    
                    page_content = response.choices[0].message.content
                    if "No relevant information on this page" not in page_content:
                        page_responses.append({
                            "page_number": page_num,
                            "content": page_content
                        })
                    
                    if st.session_state.debug_mode:
                        with st.expander(f"Debug: Page {page_num} Analysis", expanded=False):
                            st.write(f"Page {page_num} Content Length: {len(page['text'])}")
                            st.write("Relevant Content Found:" if page_responses else "No Relevant Content")
                            if page_responses:
                                st.write(page_content)
                
                if not page_responses:
                    return "No relevant information found in the document to answer this question."
                
                # Combine all relevant page responses
                combined_response = "\n\n".join([
                    f"From page {resp['page_number']}:\n{resp['content']}"
                    for resp in page_responses
                ])
                
                return combined_response

            # Process query against stored documents
            # Initialize query variable from session state
            query = st.session_state.get('current_query', None)

            if query and (pdf_texts or ppt_texts) and not sales_data:
                try:
                    responses = []
                    
                    # Process PPTs
                    for ppt_name, ppt_data in [(ppt["name"], ppt) for ppt in ppt_texts]:
                        if st.session_state.debug_mode:
                            st.write(f"Analyzing PPT: {ppt_name}")
                        
                        try:
                            # Get document from database
                            doc_data = db.get_document(ppt_name)
                            if doc_data:
                                # Use slides/pages for PPT files
                                slides_data = {
                                    "pages": doc_data["content"]["slides/pages"]  # Map slides to pages format
                                }
                                response = process_pdf_query_page_wise(query, slides_data)
                                
                                if "No relevant information found" not in response:
                                    responses.append({
                                        "source": ppt_name,
                                        "type": "ppt",
                                        "content": response
                                    })
                        except Exception as e:
                            if st.session_state.debug_mode:
                                st.error(f"Error processing PPT {ppt_name}: {str(e)}")
                            continue
                    
                    # Process PDFs
                    for pdf_name, pdf_data in [(pdf["name"], pdf) for pdf in pdf_texts]:
                        if st.session_state.debug_mode:
                            st.write(f"Analyzing PDF: {pdf_name}")
                        
                        try:
                            # Get document from database
                            doc_data = db.get_document(pdf_name)
                            if doc_data:
                                response = process_pdf_query_page_wise(query, doc_data["content"])
                                
                                if response != "No relevant information found in the document to answer this question.":
                                    responses.append({
                                        "pdf_name": pdf_name,
                                        "response": response
                                    })
                        except Exception as e:
                            if st.session_state.debug_mode:
                                st.error(f"Error processing PDF {pdf_name}: {str(e)}")
                            continue

                    # Synthesize all responses
                    if responses:
                        try:
                            synthesis_prompt = f"""
                            Based on the following information from various sources:
                            {json.dumps(responses, indent=2)}
                            
                            Provide a comprehensive answer to: {query}
                            
                            Requirements:
                            1. Synthesize all relevant information
                            2. For each piece of information, cite its source (PPT or PDF filename)
                            3. Present information in a clear, logical flow
                            4. Focus only on relevant insights
                            5. Format the response as follows:
                                - Main answer with insights
                                - Sources section listing which files contributed to the answer
                            """
                            
                            final_response = client.chat.completions.create(
                                model="gpt-4o",
                                messages=[
                                    {"role": "system", "content": "You are synthesizing information from multiple sources, providing clear answers with proper source citations."},
                                    {"role": "user", "content": synthesis_prompt}
                                ],
                                temperature=0,
                                max_tokens=1000
                            )
                            
                            final_insight = final_response.choices[0].message.content
                            st.write(final_insight)
                            
                            # Update chat history
                            for chat in st.session_state['chat_history']:
                                if chat["query"] == query:
                                    chat["response"] = final_insight
                                    break
                            auto_save_chat_history()
                        except Exception as e:
                            if st.session_state.debug_mode:
                                st.error(f"Error synthesizing responses: {str(e)}")
                            st.error("An error occurred while processing your query. Please try again.")
                    else:
                        st.write("No relevant information found in any of the uploaded documents.")

                except Exception as e:
                    if st.session_state.debug_mode:
                        st.error(f"Error in query processing: {str(e)}")
                    st.error("An error occurred while processing your query. Please try again.")
                finally:
                    st.session_state.should_query = False
            elif query and sales_data:
                # Generate drill-down questions using the chat method
                    try:
                            drill_questions = generate_drill_down_questions(temp_query, dataframes)
                            dril_results_PPT = []
                            if drill_questions['type'] == 'PPT':
                                st.write("### Drill Down Questions")
                                for dq in drill_questions['response']:
                                    st.write(f"{dq}")
                                start_time = time.time()

                                for idx, dq in enumerate(drill_questions['response'], 1):
                                    try:
                                        if 'state' in dq.lower():
                                            filtered_dataframes = [preprocess_data_for_state_query(df) for df in dataframes]
                                            temp_sales_data = SmartDatalake(filtered_dataframes, config={"llm": llm, "use_error_correction_framework": True, "data_dict": data_dict})
                                            dq_result = temp_sales_data.chat(dq)
                                        else:
                                            try:
                                                # st.write("Processing for ",dq)
                                                dq_result = sales_data.chat(dq)
                                            except e:
                                                st.error("Error: ",e)
                                            # st.write(f"Result for {idx} ",dq_result)
                                        if isinstance(dq_result,pd.DataFrame):
                                            dq_result = dq_result.to_string()
                                        else:
                                            dq_result = dq_result
                                        dril_results_PPT.append(dq_result)
                                    
                                    except Exception as e:
                                        st.error(f"Error processing drill-down question: {dq} -> {e}")
                                # st.write(drill_results)
                                total_time = time.time() - start_time
                                st.write(f"Total processing time for all questions: {total_time:.2f} seconds")

                                # Consolidate the drill-down results using the chat method
                                consolidation_prompt = consolidation_prompt.format(dril_results_PPT=dril_results_PPT,temp_query=temp_query)
                                response = client.chat.completions.create(
                                    model="gpt-4o",
                                    messages=[
                                        {"role": "system", "content": "You are an AI assistant."},
                                        {"role": "user", "content": consolidation_prompt}
                                    ],
                                    temperature=0,
                                    max_tokens=1024
                                )
                                final_insight = response.choices[0].message.content
                                st.write("### Final Consolidated Insight")
                                st.write(final_insight)

                                # Update chat history with the response
                                for chat in st.session_state['chat_history']:
                                    if chat["query"] == temp_query:
                                        chat["response"] = final_insight
                                        break
                                auto_save_chat_history()
                                st.session_state.should_query = False
                            else:
                                st.write("### Drill Down Questions")
                                for dq in drill_questions['response']:
                                    st.write(f"{dq}")

                                # Execute analysis for each drill-down question
                                drill_results = {
                                                "Overall Performance (2022-2024)": [],
                                                "Year-over-Year (YoY) Changes":[],
                                                "Regional Performance (2024)": [],
                                                "Market Penetration (2024)": [],
                                                "Consumer Analysis (2024)": [],
                                                # "Segment Performance (2024)": [],
                                                "Distribution (2024)": []
                                                }
                                keys = list(drill_results.keys())
                                start_time = time.time()
                                
                                for idx, dq in enumerate(drill_questions['response'], 1):
                                    try:
                                        if 'state' in dq.lower():
                                            filtered_dataframes = [preprocess_data_for_state_query(df) for df in dataframes]
                                            temp_sales_data = SmartDatalake(filtered_dataframes, config={"llm": llm, "use_error_correction_framework": True, "data_dict": data_dict})
                                            dq_result = temp_sales_data.chat(dq)
                                        else:
                                            try:
                                                # st.write("Processing for ",dq)
                                                dq_result = sales_data.chat(dq)
                                            except e:
                                                st.error("Error: ",e)
                                            # st.write(f"Result for {idx}, {keys[idx-1]} ",dq_result)
                                        if isinstance(dq_result,pd.DataFrame):
                                            dq_result = dq_result.to_string()
                                        else:
                                            dq_result = dq_result
                                        drill_results[keys[idx-1]].append(dq_result)
                                    
                                    except Exception as e:
                                        st.error(f"Error processing drill-down question: {dq} -> {e}")
                                # st.write(drill_results)
                                total_time = time.time() - start_time
                                st.write(f"Total processing time for all questions: {total_time:.2f} seconds")

                                # Consolidate the drill-down results using the chat method
                                consolidation_prompt = consolidatin_prompt_main.format(temp_query=temp_query,drill_results=drill_results)
                                response = client.chat.completions.create(
                                    model="gpt-4o",
                                    messages=[
                                        {"role": "system", "content": "You are an AI assistant."},
                                        {"role": "user", "content": consolidation_prompt}
                                    ],
                                    temperature=0,
                                    max_tokens=1024
                                )
                                final_insight = response.choices[0].message.content
                                st.write("### Final Consolidated Insight")
                                st.write(final_insight)

                                # Update chat history with the response
                                for chat in st.session_state['chat_history']:
                                    if chat["query"] == temp_query:
                                        chat["response"] = final_insight
                                        break
                                auto_save_chat_history()
                                st.session_state.should_query = False
                    except Exception as e:
                            st.error(f"An error occurred: {e}")
                            st.session_state.should_query = False
                    finally:
                        st.session_state.should_query = False


        # Clean up ALL checklist state for next use
        st.session_state.selected_prompts = None
        st.session_state.checkbox_states = []
        st.session_state.checkbox_option_set = set()
        st.session_state.filtered_results = None
        st.session_state.intent_type = None

    display_token_usage()

# Assume this dictionary saved in session state:
query_sql_map = st.session_state.get('query_sql_map', {})
# print(query_sql_map)
# Sidebar expanded section for queries
with st.sidebar.expander("Queries History", expanded=True):
    for i, user_query in enumerate(query_sql_map.keys()):
        if st.button(user_query, key=f"history_btn_{i}"):
            # When user clicks, load stored SQL query into session state
            st.session_state['selected_sql_for_edit'] = query_sql_map[user_query]
            st.session_state['sql_text_area_val'] = query_sql_map[user_query]   # <-- programmatically update text area
            st.session_state['selected_user_query_for_sql'] = user_query  # <-- Save the user query text

# On first run, if sql_text_area_val not set, provide a default (not strictly necessary if you want it blank)
if 'sql_text_area_val' not in st.session_state:
    st.session_state['sql_text_area_val'] = st.session_state.get('last_generated_sql', '')

# SQL Tab

with tab2:
    st.subheader('Generated SQL Query')
    # Show selected user question above SQL input if available
    selected_question = st.session_state.get('selected_user_query_for_sql', '')
    if selected_question:
        st.markdown(f"<p style='font-size:18px; font-weight:bold;'>Selected Question: {selected_question}</p>", unsafe_allow_html=True)
    sql_query = st.text_area("SQL Query:", key="sql_text_area_val", height=200)

    crud_keywords = ["INSERT", "UPDATE", "DELETE", "DROP", "ALTER", "TRUNCATE", "CREATE"]

    if st.button("Execute SQL Query"):
        # Check for CRUD keywords (case-insensitive)
        sql_query_upper = sql_query.upper()
        if any(kw in sql_query_upper for kw in crud_keywords):
            st.warning("⚠️ Access denied: CRUD operations (INSERT, UPDATE, DELETE, DROP, ALTER, TRUNCATE, CREATE) are not allowed.")
        else:
            try:
                df_result = st.session_state.engine.execute_sql(sql_query)
                df = df_result.copy()
                df = df.reset_index(drop=True)

                def format_number(x, is_year=False):
                    if pd.isna(x):
                        return ""
                    if is_year:
                        return f"{int(x)}"
                    else:
                        return f"{x:,.2f}"

                for col in df.columns:
                    if col.lower() == "year":
                        df[col] = df[col].map(lambda x: format_number(x, is_year=True))
                    elif pd.api.types.is_numeric_dtype(df[col]):
                        df[col] = df[col].map(format_number)

                if df_result is not None and not df_result.empty:
                    st.write("Query Result:")
                    st.dataframe(df, hide_index=True)
                else:
                    st.info("No data returned from query.")
            except Exception as e:
                st.error(f"Failed to execute SQL: {e}")


# Function to load latest chat history
def load_latest_chat_history():
    """Automatically load the most recent chat history file"""
    try:
        os.makedirs('chat_history', exist_ok=True)
        chat_history_files = [f for f in os.listdir('chat_history') if f.startswith('chat_history_') and f.endswith('.json')]
        if not chat_history_files:
            return
        latest_file = max(chat_history_files)
        file_path = os.path.join('chat_history', latest_file)
        if not st.session_state['chat_history']:
            with open(file_path, 'r', encoding='utf-8') as f:
                st.session_state['chat_history'] = json.load(f)
    except Exception as e:
        st.error(f"Error loading chat history: {e}")

load_latest_chat_history()

# Chat History tab
with tab3:
    st.header("Chat History")
    if st.session_state['chat_history']:
        chats_by_date = {}
        for chat in st.session_state['chat_history']:
            date = chat.get('timestamp', 'No Date').split(' ')[0]
            if date not in chats_by_date:
                chats_by_date[date] = []
            chats_by_date[date].append(chat)
        
        for date in sorted(chats_by_date.keys(), reverse=True):
            st.subheader(date)
            for chat in chats_by_date[date]:
                with st.expander(f"Q: {chat['query'][:100]}..."):
                    st.write("**Question:**")
                    st.write(chat['query'])
                    st.write("**Answer:**")
                    if chat['response']:
                        st.write(chat['response'])
                    else:
                        st.write("*No response recorded*")
                    if 'feedback' in chat:
                        st.write("**Feedback:**", chat['feedback'])
                    if 'timestamp' in chat:
                        st.write("*Timestamp:*", chat['timestamp'])
    else:
        st.info("No chat history available. Start a conversation to create history.")

# Function to handle feedback (unchanged)
def handle_feedback(chat_idx, feedback):
    st.session_state['chat_history'][chat_idx]['feedback'] = feedback

# if os.name == 'nt':  # Windows
#     pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'